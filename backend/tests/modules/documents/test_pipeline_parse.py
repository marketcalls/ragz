import io

import pytest
from docx import Document as DocxBuilder

from ragz.modules.documents import pipeline
from ragz.modules.documents.pipeline import IngestFailure, PageBlock, needs_ocr, parse_bytes


def build_docx() -> bytes:
    d = DocxBuilder()
    d.add_heading("Flux Capacitor Manual", level=1)
    d.add_paragraph("The flux capacitor requires 1.21 gigawatts of power.")
    d.add_heading("Billing", level=1)
    d.add_paragraph("Invoice 0231 covers the plutonium delivery.")
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_parse_docx_extracts_blocks_with_kinds() -> None:
    blocks = parse_bytes(build_docx(), "manual.docx")
    texts = " ".join(b.text for b in blocks)
    assert "1.21 gigawatts" in texts and "Invoice 0231" in texts
    assert any(b.kind == "heading" for b in blocks)
    assert all(b.page >= 1 for b in blocks)


def test_parse_txt_fast_path() -> None:
    blocks = parse_bytes(b"para one\n\npara two", "notes.txt")
    assert [b.text for b in blocks] == ["para one", "para two"]
    assert all(b.kind == "text" and b.page == 1 for b in blocks)


def test_empty_file_fails_with_reason() -> None:
    with pytest.raises(IngestFailure, match="empty"):
        parse_bytes(b"", "empty.txt")


def test_unsupported_format_fails_with_reason() -> None:
    with pytest.raises(IngestFailure):
        parse_bytes(b"\x00\x01garbage", "weird.xyz")


def _pptx_bytes() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    for i in (1, 2):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Safety Briefing Part {i}"
        slide.placeholders[1].text = (
            f"Slide {i} body: always wear certified PPE inside zone {i}."
        )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_parse_pptx_slides_become_pages() -> None:
    blocks = parse_bytes(_pptx_bytes(), "briefing.pptx")
    assert {b.page for b in blocks} == {1, 2}  # slide number = page number
    joined = " ".join(b.text for b in blocks)
    assert "certified PPE" in joined
    # Slide titles surface as headings (feeds Task 3's section trail).
    assert any(b.kind == "heading" and "Safety Briefing" in b.text for b in blocks)


def test_headings_carry_level() -> None:
    # .txt path emits no headings; use the pptx fixture from Task 1 for a real doc.
    blocks = parse_bytes(_pptx_bytes(), "briefing.pptx")
    heading = next(b for b in blocks if b.kind == "heading")
    assert heading.level is not None
    text_block = next(b for b in blocks if b.kind == "text")
    assert text_block.level is None


def test_needs_ocr_low_density() -> None:
    sparse = [PageBlock(page=3, text="p3", kind="text")]  # 2 chars over 3 pages
    assert needs_ocr(sparse, min_chars_per_page=200)
    dense = [PageBlock(page=1, text="x" * 500, kind="text")]
    assert not needs_ocr(dense, min_chars_per_page=200)
    assert needs_ocr([], min_chars_per_page=200)  # zero blocks = scanned


def test_parse_pdf_low_text_triggers_ocr_pass(monkeypatch) -> None:  # type: ignore[no-untyped-def]
    calls: list[bool] = []

    def fake_convert(tmp_path, suffix, *, ocr):  # type: ignore[no-untyped-def]
        calls.append(ocr)
        if not ocr:
            return [PageBlock(page=2, text="x", kind="text")]  # sparse first pass
        return [PageBlock(page=1, text="EMERGENCY ASSEMBLY POINT", kind="text")]

    monkeypatch.setattr(pipeline, "_convert_blocks", fake_convert)
    blocks = parse_bytes(b"%PDF-fake", "scan.pdf", ocr_enabled=True)
    assert calls == [False, True]
    assert blocks[0].text == "EMERGENCY ASSEMBLY POINT"


def test_parse_pdf_kill_switch_skips_ocr(monkeypatch) -> None:  # type: ignore[no-untyped-def]
    calls: list[bool] = []

    def fake_convert(tmp_path, suffix, *, ocr):  # type: ignore[no-untyped-def]
        calls.append(ocr)
        return [PageBlock(page=2, text="x", kind="text")]

    monkeypatch.setattr(pipeline, "_convert_blocks", fake_convert)
    blocks = parse_bytes(b"%PDF-fake", "scan.pdf", ocr_enabled=False)
    assert calls == [False]
    assert blocks  # sparse result returned as-is rather than OCR'd


def test_parse_non_pdf_never_ocrs(monkeypatch) -> None:  # type: ignore[no-untyped-def]
    calls: list[bool] = []

    def fake_convert(tmp_path, suffix, *, ocr):  # type: ignore[no-untyped-def]
        calls.append(ocr)
        return [PageBlock(page=1, text="x", kind="text")]

    monkeypatch.setattr(pipeline, "_convert_blocks", fake_convert)
    parse_bytes(b"data", "tiny.docx", ocr_enabled=True)
    assert calls == [False]
